import os
import re
import math
import pypdf
import docx
import pptx
from typing import List, Dict, Any, Optional
import json
import faiss
import numpy as np

# ── Generative AI & Environment ───────────────────────────────────────
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    # Use gemini-2.0-flash for higher free-tier quota (1500 RPD vs 20 RPD for 3.5-flash)
    llm_model = genai.GenerativeModel("gemini-2.0-flash")
else:
    llm_model = None

# ── Lazy-loaded globals ────────────────────────────────────────────────
_embedding_model = None
_faiss_index = None
_faiss_metadata = [] # List of dicts matching FAISS indices

def _get_model():
    global _embedding_model
    if _embedding_model is None:
        from sentence_transformers import SentenceTransformer
        _embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    return _embedding_model

def _init_faiss():
    global _faiss_index, _faiss_metadata
    if _faiss_index is None:
        os.makedirs("faiss_db", exist_ok=True)
        # Initialize flat L2 index
        _faiss_index = faiss.IndexFlatL2(384) # all-MiniLM-L6-v2 dim
        
        # Load from disk if exists
        if os.path.exists("faiss_db/index.faiss"):
            _faiss_index = faiss.read_index("faiss_db/index.faiss")
        if os.path.exists("faiss_db/metadata.json"):
            with open("faiss_db/metadata.json", "r") as f:
                _faiss_metadata = json.load(f)

def _save_faiss():
    global _faiss_index, _faiss_metadata
    os.makedirs("faiss_db", exist_ok=True)
    if _faiss_index:
        faiss.write_index(_faiss_index, "faiss_db/index.faiss")
        with open("faiss_db/metadata.json", "w") as f:
            json.dump(_faiss_metadata, f)

# ── Text Extraction ────────────────────────────────────────────────────
def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    t = page.extract_text()
                    if t: text += str(t).strip() + "\n"
        elif ext == ".docx":
            with open(file_path, "rb") as f:
                doc = docx.Document(f)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in [".ppt", ".pptx"]:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        print(f"Text extraction error ({file_path}): {e}")
    return text.strip()

# ── Extractive Summarization (TF-IDF) ─────────────────────────────────
def summarize_text(text: str, num_sentences: int = 5) -> str:
    # Use Gemini if available
    if llm_model:
        try:
            prompt = f"Provide a concise summary of the following text in bullet points:\n\n{text[:10000]}"
            resp = llm_model.generate_content(prompt)
            return resp.text
        except:
            pass

    # Fallback TF-IDF scoring
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 30]
    if len(sentences) <= num_sentences:
        return " ".join(sentences)

    words_per = [re.findall(r'\w+', s.lower()) for s in sentences]
    tf = [{w: ws.count(w)/len(ws) for w in set(ws)} for ws in words_per]
    doc_freq = {}
    for ws in words_per:
        for w in set(ws):
            doc_freq[w] = doc_freq.get(w, 0) + 1
    N = len(sentences)
    scores = []
    for i, ws in enumerate(words_per):
        score = sum(tf[i].get(w, 0) * math.log((N + 1) / (doc_freq.get(w, 1) + 1))
                    for w in ws) / (len(ws) or 1)
        scores.append((score, i))
    top = sorted(scores, reverse=True)[:num_sentences]
    top_idx = sorted(idx for _, idx in top)
    return " ".join(sentences[i] for i in top_idx)

# ── Auto-Categorization ────────────────────────────────────────────────
def suggest_category(text: str, existing_categories: List[str]) -> str:
    if llm_model:
        try:
            prompt = f"Categorize the following document excerpt into exactly one of these categories: {', '.join(existing_categories)}. Respond with only the category name.\n\nExcerpt:\n{text[:3000]}"
            resp = llm_model.generate_content(prompt)
            cat = resp.text.strip()
            if cat in existing_categories:
                return cat
        except:
            pass

    if not text: return "Knowledge Articles / FAQs"
    text_lower = text.lower()
    taxonomy = {
        "Academic Resources": ["course", "syllabus", "exam", "question paper", "assignment", "lab", "tutorial", "study", "lecture", "notes", "mathematics", "physics", "science", "chemistry", "algorithms"],
        "Research & Publications": ["research", "paper", "thesis", "dissertation", "conference", "journal", "dataset", "dataset", "publication"],
        "Administrative Resources": ["regulation", "policy", "circular", "form", "guideline", "notice", "administration", "rules", "leave", "attendance"],
        "Student Services": ["scholarship", "internship", "placement", "career", "counseling", "hostel", "alumni"],
        "Events & Announcements": ["workshop", "seminar", "event", "announcement", "schedule", "cultural", "fest"],
        "Technical Documentation": ["software", "equipment", "programming", "system", "documentation", "code", "architecture", "docker", "kubernetes", "cloud", "api"],
        "Knowledge Articles / FAQs": ["faq", "procedure", "safety", "help", "guide", "how to", "instructions"],
        "Multimedia Learning": ["video", "podcast", "infographic", "recorded", "audio"],
        "Institutional Reports": ["annual report", "accreditation", "statistics", "performance", "budget", "finance"],
        "External Resources": ["mooc", "external", "portal", "coursera", "udemy"]
    }
    best_category = "Academic Resources"
    highest_score = 0
    for category, keywords in taxonomy.items():
        score = sum(text_lower.count(kw) * (1.5 if len(kw) > 6 else 1.0) for kw in keywords)
        if score > highest_score:
            highest_score, best_category = score, category
    return best_category if highest_score >= 2 else "Knowledge Articles / FAQs"

def extract_tags(text: str) -> str:
    if llm_model:
        try:
            prompt = f"Extract exactly 5 descriptive keywords/tags from this text. Output as comma separated list:\n\n{text[:3000]}"
            resp = llm_model.generate_content(prompt)
            return resp.text.strip()
        except:
            pass
            
    if not text: return "document"
    text_lower = text[:5000].lower()
    words = re.findall(r'\b[a-zA-Z]{5,15}\b', text_lower)
    STOP_WORDS = {"which", "their", "there", "about", "could", "would", "these", "those", "other", "where", "after", "first", "years", "often", "being", "under", "because", "while", "between", "through", "should", "using", "however", "therefore", "always", "never", "might", "every", "something", "anything", "nothing", "someone", "anyone", "without", "within", "during", "before", "really", "almost", "already", "although", "always", "another", "example", "number", "system", "time", "what", "when", "why", "this", "that", "with", "from", "have", "state", "shows", "given", "taken", "made", "used", "paper", "document", "study", "analysis", "report"}
    
    def score_word(w: str) -> int:
        if w in STOP_WORDS: return -100
        score = len(w)
        if w.endswith(('tion', 'logy', 'ics', 'ment', 'phy', 'try', 'ing')): score += 20
        score += text_lower.count(w) * 2
        return score
        
    valid_words = list(set(words))
    valid_words.sort(key=score_word, reverse=True)
    tags = [w for w in valid_words if score_word(w) > 5][:5]
    return ", ".join(tags) if tags else "document"

# ── AI Indexing (FAISS) ────────────────────────────────────────────────
def process_and_index_document(document_id: int, title: str, file_path: str, category: str) -> bool:
    text = extract_text_from_file(file_path)
    if not text:
        return False
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_text(text)

        _init_faiss()
        model = _get_model()
        embeddings = model.encode(chunks)
        
        # Add to FAISS
        faiss.normalize_L2(embeddings)
        _faiss_index.add(embeddings)
        
        # Add to metadata
        global _faiss_metadata
        for i, chunk in enumerate(chunks):
            _faiss_metadata.append({
                "document_id": document_id,
                "title": title,
                "category": category,
                "chunk_index": i,
                "content": chunk
            })
            
        _save_faiss()
        return True
    except Exception as e:
        print(f"Indexing error: {e}")
        return False

def remove_document_from_index(document_id: int) -> bool:
    try:
        _init_faiss()
        global _faiss_metadata
        
        indices_to_remove = [i for i, m in enumerate(_faiss_metadata) if m.get("document_id") == document_id]
        if not indices_to_remove: return True
        
        # FAISS index removing ID is tricky in IndexFlatL2, we have to rebuild it
        _faiss_metadata = [m for i, m in enumerate(_faiss_metadata) if i not in indices_to_remove]
        
        global _faiss_index
        _faiss_index = faiss.IndexFlatL2(384)
        
        if _faiss_metadata:
            chunks = [m["content"] for m in _faiss_metadata]
            model = _get_model()
            embeddings = model.encode(chunks)
            faiss.normalize_L2(embeddings)
            _faiss_index.add(embeddings)
            
        _save_faiss()
        return True
    except Exception as e:
        print(f"Index deletion error: {e}")
        return False

# ── Duplicate Detection ────────────────────────────────────────────────
def check_duplicate(file_path: str, threshold: float = 0.92) -> Optional[Dict]:
    try:
        text = extract_text_from_file(file_path)
        if not text: return None
        snippet = text[:2000]
        
        _init_faiss()
        if _faiss_index is None or _faiss_index.ntotal == 0:
            return None
            
        model = _get_model()
        emb = model.encode([snippet])
        faiss.normalize_L2(emb)
        
        distances, indices = _faiss_index.search(emb, 1)
        if len(distances) > 0 and len(distances[0]) > 0:
            dist = distances[0][0]
            # Convert L2 distance of normalized vectors to Cosine Similarity
            similarity = 1 - (dist / 2)
            if similarity >= threshold:
                idx = indices[0][0]
                meta = _faiss_metadata[idx]
                return {"document_id": meta.get("document_id"), "title": meta.get("title"), "similarity": round(similarity * 100, 1)}
    except Exception as e:
        print(f"Duplicate check error: {e}")
    return None

# ── Semantic Search ────────────────────────────────────────────────────
def semantic_search(query: str, top_k: int = 5, category: str = None, document_id: int = None) -> List[Dict[str, Any]]:
    try:
        _init_faiss()
        if _faiss_index is None or _faiss_index.ntotal == 0:
            return []
            
        model = _get_model()
        emb = model.encode([query])
        faiss.normalize_L2(emb)
        
        # Over-fetch and then filter by category
        search_k = min(top_k * 5, _faiss_index.ntotal)
        distances, indices = _faiss_index.search(emb, search_k)
        
        out = []
        for i, idx in enumerate(indices[0]):
            if idx == -1: continue
            dist = distances[0][i]
            if dist > 1.0: # Filter out unrelated
                continue
                
            meta = _faiss_metadata[idx]
            if category and meta.get("category") != category:
                continue
            if document_id is not None and meta.get("document_id") != document_id:
                continue
                
            out.append({
                "content": meta["content"],
                "metadata": meta,
                "distance": float(dist)
            })
            
            if len(out) >= top_k:
                break
        return out
    except Exception as e:
        print(f"Search error: {e}")
        return []

# ── RAG Answer ────────────────────────────────────────────────────────
def generate_rag_answer(query: str, context: str) -> str:
    if not context.strip():
        return "I couldn't find relevant information in the knowledge base for your query. Please ensure documents have been uploaded and indexed."
        
    if llm_model:
        try:
            prompt = f"Answer the user query based ONLY on the provided context. If the context does not contain the answer, say 'The uploaded documents do not contain this information.'\n\nContext:\n{context}\n\nQuery: {query}"
            resp = llm_model.generate_content(prompt)
            return resp.text.strip()
        except:
            pass
            
    # Fallback if Gemini quota is exceeded or not configured
    return f"I couldn't use the AI model right now to formulate an answer (API quota exceeded). Here is the raw excerpt from the knowledge base that matches your query:\n\n\"{context[:500]}...\""

# ── Question Generator ─────────────────────────────────────────────────
def generate_questions(text: str, num_questions: int = 5) -> List[Dict]:
    if llm_model:
        import time
        for attempt in range(3):
            try:
                prompt = f"""Generate exactly {num_questions} Multiple Choice Questions (MCQs) from this text.
Return ONLY a JSON array (not an object). Each element must have these fields:
- "question": the question text
- "type": "mcq"
- "options": array of exactly 4 option strings
- "answer": the correct option (must exactly match one of the options)

Text:
{text[:12000]}"""
                resp = llm_model.generate_content(prompt, generation_config={"response_mime_type": "application/json"})
                data = json.loads(resp.text)
                
                # Handle both list and dict responses from Gemini
                if isinstance(data, list):
                    questions = data
                elif isinstance(data, dict):
                    questions = data.get("questions", data.get("mcqs", []))
                    if not questions:
                        # Try to find any list value in the dict
                        for v in data.values():
                            if isinstance(v, list) and len(v) > 0:
                                questions = v
                                break
                else:
                    questions = []
                
                # Validate structure
                valid_questions = []
                for q in questions:
                    if isinstance(q, dict) and "question" in q and "options" in q and "answer" in q:
                        valid_questions.append(q)
                
                if valid_questions:
                    return valid_questions
                    
            except Exception as e:
                err_str = str(e)
                print(f"Quiz generation attempt {attempt+1} error: {err_str}")
                if "429" in err_str or "ResourceExhausted" in err_str:
                    if attempt < 2:
                        time.sleep(5 * (attempt + 1))  # Backoff: 5s, 10s
                        continue
                break  # Don't retry for non-rate-limit errors

    # Fallback to comprehension-based question generation (works without API)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 50 and len(s.strip()) < 400]
    questions = []
    
    if not sentences:
        return questions

    import random
    random.shuffle(sentences)
    
    used_sents = set()
    for sent in sentences:
        if len(questions) >= num_questions:
            break
        if sent in used_sents:
            continue
        used_sents.add(sent)
        
        # Create a "Which statement is correct?" style question
        # The correct answer is the actual sentence content
        # Generate plausible-looking wrong answers from other sentences
        
        # Shorten sentence for the answer option if too long
        answer_text = sent if len(sent) <= 150 else sent[:147] + "..."
        
        # Pick other sentences as distractors (they're from the same doc so they look plausible)
        other_sents = [s for s in sentences if s != sent and s not in used_sents]
        random.shuffle(other_sents)
        distractors = []
        for d in other_sents[:3]:
            d_text = d if len(d) <= 150 else d[:147] + "..."
            distractors.append(d_text)
        
        # If not enough distractors, create variations
        while len(distractors) < 3:
            distractors.append(f"None of the above statements are correct.")
            if len(distractors) < 3:
                distractors.append(f"This topic is not covered in the document.")
            if len(distractors) < 3:
                distractors.append(f"The document does not provide this information.")
        
        options = [answer_text] + distractors[:3]
        random.shuffle(options)
        
        # Create the question text based on key topic words
        topic_words = [w for w in re.findall(r'\b[A-Z][a-z]{3,}\b', sent)]
        if topic_words:
            topic = topic_words[0]
            q_text = f"According to the document, which of the following statements about {topic} is correct?"
        else:
            q_text = "According to the document, which of the following statements is correct?"
        
        questions.append({
            "question": q_text,
            "type": "mcq",
            "options": options,
            "answer": answer_text
        })
    
    return questions

# ── Knowledge Graph ────────────────────────────────────────────────────
def build_knowledge_graph(categories: list, documents: list) -> Dict:
    nodes, links, node_ids = [], [], {}
    for cat in categories:
        node_ids[f"cat_{cat.id}"] = len(nodes)
        nodes.append({"id": f"cat_{cat.id}", "label": cat.name, "type": "category", "size": 20})
    for doc in documents:
        nodes.append({"id": f"doc_{doc.id}", "label": doc.title[:30], "type": "document", "size": 10})
        if doc.category_id and f"cat_{doc.category_id}" in node_ids:
            links.append({"source": f"cat_{doc.category_id}", "target": f"doc_{doc.id}"})
    return {"nodes": nodes, "links": links}

# ── Flashcards Generator ───────────────────────────────────────────────
def generate_flashcards(text: str) -> list:
    if llm_model:
        import time
        for attempt in range(3):
            try:
                prompt = f"""Extract exactly 5 key terms and their clear definitions from this text.
Return ONLY a JSON array (not an object). Each element must have:
- "term": the key term or concept
- "definition": a clear, concise definition

Text: {text[:10000]}"""
                resp = llm_model.generate_content(prompt, generation_config={"response_mime_type": "application/json"})
                data = json.loads(resp.text)
                
                # Handle both list and dict responses
                if isinstance(data, list):
                    flashcards = data
                elif isinstance(data, dict):
                    flashcards = data.get("flashcards", [])
                    if not flashcards:
                        for v in data.values():
                            if isinstance(v, list) and len(v) > 0:
                                flashcards = v
                                break
                else:
                    flashcards = []
                
                # Validate
                valid = [fc for fc in flashcards if isinstance(fc, dict) and "term" in fc and "definition" in fc]
                if valid:
                    return valid
                    
            except Exception as e:
                err_str = str(e)
                print(f"Flashcard attempt {attempt+1} error: {err_str}")
                if "429" in err_str or "ResourceExhausted" in err_str:
                    if attempt < 2:
                        time.sleep(5 * (attempt + 1))
                        continue
                break
    
    # Smart text-based fallback: find definitional sentences from the document
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 40 and len(s.strip()) < 400]
    
    if not sentences:
        return [{"term": "Document Overview", "definition": text[:300] + "..." if len(text) > 300 else text}]
    
    # Prioritize sentences that look like definitions ("X is ...", "X refers to ...", "X is defined as ...")
    def_patterns = [
        re.compile(r'^([A-Z][A-Za-z\s]{2,30})\s+(?:is|are|refers?\s+to|means?|can\s+be\s+defined\s+as|is\s+defined\s+as|is\s+a|is\s+an|is\s+the)\s+(.{20,})', re.IGNORECASE),
        re.compile(r'^(?:The|A|An)\s+([A-Za-z\s]{3,30})\s+(?:is|are|refers?\s+to|means?)\s+(.{20,})', re.IGNORECASE),
    ]
    
    flashcards = []
    used_terms = set()
    
    # First pass: find definitional sentences
    for sent in sentences:
        if len(flashcards) >= 5:
            break
        for pattern in def_patterns:
            match = pattern.match(sent)
            if match:
                term = match.group(1).strip().rstrip(',')
                if len(term) > 3 and term.lower() not in used_terms:
                    used_terms.add(term.lower())
                    flashcards.append({"term": term, "definition": sent})
                    break
    
    # Second pass: use key sentences with capitalized terms if we need more
    if len(flashcards) < 5:
        import random
        remaining = [s for s in sentences if s not in [f["definition"] for f in flashcards]]
        random.shuffle(remaining)
        
        for sent in remaining:
            if len(flashcards) >= 5:
                break
            # Extract capitalized terms (likely proper nouns / key concepts)
            caps = re.findall(r'\b([A-Z][a-z]{3,}(?:\s+[A-Z]?[a-z]{3,})?)\b', sent)
            # Filter out sentence starters
            words = sent.split()
            first_word = words[0] if words else ""
            caps = [c for c in caps if c != first_word and c.lower() not in used_terms]
            
            if caps:
                term = caps[0]
                used_terms.add(term.lower())
                flashcards.append({"term": term, "definition": sent})
    
    # Final fallback: just use numbered key points
    if len(flashcards) < 3:
        for i, sent in enumerate(sentences[:5]):
            if len(flashcards) >= 5:
                break
            if sent not in [f["definition"] for f in flashcards]:
                flashcards.append({"term": f"Key Point {len(flashcards) + 1}", "definition": sent})
    
    return flashcards if flashcards else [{"term": "Study Note", "definition": text[:300]}]

